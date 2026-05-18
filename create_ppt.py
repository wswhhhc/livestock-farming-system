#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
家畜养殖信息采集及收益估计系统 - 答辩PPT生成脚本
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ===== 配置 =====
ASSETS_DIR = r"C:\Users\86193\.cursor\projects\d-AI-Demo-jiachuyangzhi\assets"
OUTPUT_FILE = r"d:\AI Demo\jiachuyangzhi\家畜养殖系统答辩.pptx"

# 绿色系配色方案（清新自然风）
COLORS = {
    'primary_green': RGBColor(46, 125, 50),      # 主绿色 #2E7D32
    'light_green': RGBColor(129, 199, 132),      # 浅绿 #81C784
    'dark_green': RGBColor(27, 94, 32),           # 深绿 #1B5E20
    'accent_green': RGBColor(76, 175, 80),       # 强调绿 #4CAF50
    'bg_light': RGBColor(250, 252, 250),          # 背景浅绿白
    'text_dark': RGBColor(33, 37, 41),            # 深灰文字
    'text_gray': RGBColor(108, 117, 125),         # 灰色文字
    'white': RGBColor(255, 255, 255),
    'cream': RGBColor(245, 245, 220),             # 米黄色
}

# 8张截图（按顺序）
SCREENSHOTS = [
    ("c__Users_86193_AppData_Roaming_Cursor_User_workspaceStorage_empty-window_images_image-950fa446-50c0-4426-a172-6ddabf134671.png", "仪表盘"),
    ("c__Users_86193_AppData_Roaming_Cursor_User_workspaceStorage_empty-window_images_image-7fa4fdf3-16e3-48a4-998f-6ebc2cdf170b.png", "家畜种类管理"),
    ("c__Users_86193_AppData_Roaming_Cursor_User_workspaceStorage_empty-window_images_image-e4d2a90d-d0a2-4313-8be7-08469e80c24e.png", "养殖场地管理"),
    ("c__Users_86193_AppData_Roaming_Cursor_User_workspaceStorage_empty-window_images_image-89aed0bd-64c3-4a58-a976-98511828f7fc.png", "存栏列表"),
    ("c__Users_86193_AppData_Roaming_Cursor_User_workspaceStorage_empty-window_images_image-89448bbd-d195-46f4-a8f1-5ebb794e22d0.png", "成本管理"),
    ("c__Users_86193_AppData_Roaming_Cursor_User_workspaceStorage_empty-window_images_image-6c957fbc-84e4-4f4f-8911-7bb297c71498.png", "价格规则管理"),
    ("c__Users_86193_AppData_Roaming_Cursor_User_workspaceStorage_empty-window_images_image-5210ac07-619b-4525-be86-34ac0d55e65c.png", "收益预估"),
    ("c__Users_86193_AppData_Roaming_Cursor_User_workspaceStorage_empty-window_images_image-6d6be75c-d1f6-4281-bab9-fa89cb051b53.png", "养殖建议"),
]

def add_title_shape(slide, text, left, top, width, height, font_size=44, bold=True, color=None):
    """添加标题文本框"""
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color or COLORS['text_dark']
    p.alignment = PP_ALIGN.LEFT
    return shape

def add_body_text(slide, text, left, top, width, height, font_size=18, color=None, align=PP_ALIGN.LEFT):
    """添加正文文本框"""
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color or COLORS['text_gray']
    p.alignment = align
    return shape

def add_bullet_points(slide, items, left, top, width, height, font_size=16, color=None):
    """添加带项目符号的列表"""
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color or COLORS['text_dark']
        p.space_after = Pt(8)
        p.level = 0
    return shape

def add_colored_bar(slide, left, top, width, height, color):
    """添加彩色装饰条"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def set_slide_bg_color(slide, color):
    """设置幻灯片背景色"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_image_to_slide(slide, img_path, left, top, width, height):
    """添加图片到幻灯片"""
    try:
        if os.path.exists(img_path):
            slide.shapes.add_picture(img_path, left, top, width, height)
            return True
    except Exception as e:
        print(f"添加图片失败: {img_path}, 错误: {e}")
    return False

def create_presentation():
    """创建完整的PPT"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 定义常用尺寸
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    margin = Inches(0.5)
    
    # ========== 第1页：封面 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    set_slide_bg_color(slide, COLORS['bg_light'])
    
    # 顶部绿色装饰条
    add_colored_bar(slide, 0, 0, slide_width, Inches(0.15), COLORS['primary_green'])
    
    # 主标题
    add_title_shape(slide, "家畜养殖信息采集", margin, Inches(2), Inches(12), Inches(0.8), 
                    font_size=54, color=COLORS['dark_green'])
    add_title_shape(slide, "及收益估计系统", margin, Inches(2.8), Inches(12), Inches(0.8), 
                    font_size=54, color=COLORS['dark_green'])
    
    # 副标题
    add_body_text(slide, "—— 基于 Vue3 + Spring Boot 的数字化养殖管理平台", 
                  margin, Inches(3.8), Inches(12), Inches(0.5), 
                  font_size=22, color=COLORS['text_gray'])
    
    # 底部信息
    add_body_text(slide, "课程答辩演示", margin, Inches(5.5), Inches(6), Inches(0.4), 
                  font_size=18, color=COLORS['text_gray'])
    
    # 右侧装饰圆形
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10), Inches(4), Inches(3), Inches(3))
    circle.fill.solid()
    circle.fill.fore_color.rgb = COLORS['light_green']
    circle.fill.fore_color.brightness = 0.3
    circle.line.fill.background()
    
    # 底部绿色装饰条
    add_colored_bar(slide, 0, Inches(7.35), slide_width, Inches(0.15), COLORS['primary_green'])
    
    # ========== 第2页：目录 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg_color(slide, COLORS['white'])
    
    add_title_shape(slide, "目录", margin, Inches(0.5), Inches(3), Inches(0.6), 
                    font_size=40, color=COLORS['dark_green'])
    add_colored_bar(slide, margin, Inches(1.15), Inches(1.5), Inches(0.05), COLORS['accent_green'])
    
    # 目录项
    sections = [
        ("01", "项目背景与意义", "传统养殖管理的痛点与数字化需求"),
        ("02", "系统架构设计", "前后端分离架构与技术选型"),
        ("03", "核心功能演示", "8大功能模块详细展示"),
        ("04", "应用价值总结", "系统特色与创新点"),
    ]
    
    y_pos = Inches(1.8)
    for num, title, desc in sections:
        # 编号
        num_shape = slide.shapes.add_textbox(Inches(0.8), y_pos, Inches(0.8), Inches(0.6))
        num_tf = num_shape.text_frame
        num_p = num_tf.paragraphs[0]
        num_p.text = num
        num_p.font.size = Pt(32)
        num_p.font.bold = True
        num_p.font.color.rgb = COLORS['accent_green']
        
        # 标题
        add_title_shape(slide, title, Inches(1.8), y_pos, Inches(6), Inches(0.5), 
                        font_size=24, color=COLORS['text_dark'])
        
        # 描述
        add_body_text(slide, desc, Inches(1.8), y_pos + Inches(0.5), Inches(8), Inches(0.4), 
                      font_size=14, color=COLORS['text_gray'])
        
        y_pos += Inches(1.3)
    
    # ========== 第3页：团队介绍 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg_color(slide, COLORS['bg_light'])
    
    add_title_shape(slide, "团队介绍", margin, Inches(0.4), Inches(4), Inches(0.5), 
                    font_size=36, color=COLORS['dark_green'])
    add_colored_bar(slide, margin, Inches(0.95), Inches(1.5), Inches(0.05), COLORS['accent_green'])
    
    # 组长
    leader_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                                         Inches(0.8), Inches(1.4), Inches(5.8), Inches(2.4))
    leader_card.fill.solid()
    leader_card.fill.fore_color.rgb = COLORS['primary_green']
    leader_card.line.fill.background()
    
    add_body_text(slide, "组长", Inches(1), Inches(1.6), Inches(1.5), Inches(0.4), 
                  font_size=14, color=COLORS['light_green'])
    add_title_shape(slide, "王士维", Inches(1), Inches(2), Inches(3), Inches(0.6), 
                    font_size=32, color=COLORS['white'])
    add_body_text(slide, "• 项目整体架构设计\n• 前后端核心功能开发\n• 数据库设计与实现", 
                  Inches(1), Inches(2.7), Inches(5.2), Inches(0.9), 
                  font_size=15, color=COLORS['white'])
    
    # 组员1 - 徐嘉豪
    member1_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                                          Inches(6.8), Inches(1.4), Inches(5.8), Inches(2.4))
    member1_card.fill.solid()
    member1_card.fill.fore_color.rgb = COLORS['accent_green']
    member1_card.line.fill.background()
    
    add_body_text(slide, "组员", Inches(7), Inches(1.6), Inches(1.5), Inches(0.4), 
                  font_size=14, color=COLORS['white'])
    add_title_shape(slide, "徐嘉豪", Inches(7), Inches(2), Inches(3), Inches(0.6), 
                    font_size=32, color=COLORS['white'])
    add_body_text(slide, "• PPT设计与制作\n• 答辩材料整理\n• 演示流程规划", 
                  Inches(7), Inches(2.7), Inches(5.2), Inches(0.9), 
                  font_size=15, color=COLORS['white'])
    
    # 组员2 - 张永琪
    member2_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                                          Inches(0.8), Inches(4.1), Inches(5.8), Inches(2.4))
    member2_card.fill.solid()
    member2_card.fill.fore_color.rgb = COLORS['white']
    member2_card.line.color.rgb = COLORS['primary_green']
    
    add_body_text(slide, "组员", Inches(1), Inches(4.3), Inches(1.5), Inches(0.4), 
                  font_size=14, color=COLORS['text_gray'])
    add_title_shape(slide, "张永琪", Inches(1), Inches(4.7), Inches(3), Inches(0.6), 
                    font_size=32, color=COLORS['primary_green'])
    add_body_text(slide, "• 文档整理与编写\n• 后端部分功能开发\n• API接口测试", 
                  Inches(1), Inches(5.4), Inches(5.2), Inches(0.9), 
                  font_size=15, color=COLORS['text_dark'])
    
    # 组员3 - 王耀
    member3_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                                          Inches(6.8), Inches(4.1), Inches(5.8), Inches(2.4))
    member3_card.fill.solid()
    member3_card.fill.fore_color.rgb = COLORS['white']
    member3_card.line.color.rgb = COLORS['accent_green']
    
    add_body_text(slide, "组员", Inches(7), Inches(4.3), Inches(1.5), Inches(0.4), 
                  font_size=14, color=COLORS['text_gray'])
    add_title_shape(slide, "王耀", Inches(7), Inches(4.7), Inches(3), Inches(0.6), 
                    font_size=32, color=COLORS['primary_green'])
    add_body_text(slide, "• 资料查找与调研\n• 前端部分页面开发\n• UI组件调试", 
                  Inches(7), Inches(5.4), Inches(5.2), Inches(0.9), 
                  font_size=15, color=COLORS['text_dark'])
    
    # ========== 第4页：项目背景 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg_color(slide, COLORS['bg_light'])
    
    add_title_shape(slide, "项目背景", margin, Inches(0.5), Inches(4), Inches(0.6), 
                    font_size=36, color=COLORS['dark_green'])
    add_colored_bar(slide, margin, Inches(1.1), Inches(1.2), Inches(0.05), COLORS['accent_green'])
    
    # 左侧：行业痛点
    add_title_shape(slide, "行业痛点", Inches(0.8), Inches(1.6), Inches(5), Inches(0.5), 
                    font_size=24, color=COLORS['primary_green'])
    
    pain_points = [
        "手工记录效率低，数据易丢失",
        "成本核算困难，收益难以预估",
        "缺乏数据支撑，决策凭经验",
        "养殖过程难以全程追溯",
    ]
    add_bullet_points(slide, pain_points, Inches(0.8), Inches(2.2), Inches(5), Inches(3), 
                      font_size=16, color=COLORS['text_dark'])
    
    # 右侧：解决方案
    add_title_shape(slide, "解决方案", Inches(7), Inches(1.6), Inches(5), Inches(0.5), 
                    font_size=24, color=COLORS['primary_green'])
    
    solutions = [
        "数字化信息采集与存储",
        "智能化成本与收益计算",
        "数据可视化辅助决策",
        "全生命周期追溯管理",
    ]
    add_bullet_points(slide, solutions, Inches(7), Inches(2.2), Inches(5), Inches(3), 
                      font_size=16, color=COLORS['text_dark'])
    
    # 底部核心价值
    add_colored_bar(slide, Inches(0.8), Inches(5.5), Inches(11.7), Inches(1.5), COLORS['primary_green'])
    add_body_text(slide, "核心价值：提升养殖管理效率，实现科学决策，助力智慧农业",
                  Inches(1), Inches(5.9), Inches(11), Inches(0.6), 
                  font_size=20, color=COLORS['white'], align=PP_ALIGN.CENTER)
    
    # ========== 第5页：技术架构概览 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg_color(slide, COLORS['white'])
    
    add_title_shape(slide, "技术架构", margin, Inches(0.4), Inches(4), Inches(0.5), 
                    font_size=36, color=COLORS['dark_green'])
    add_colored_bar(slide, margin, Inches(0.95), Inches(1.2), Inches(0.05), COLORS['accent_green'])
    
    # 架构图 - 三层架构
    # 前端层
    frontend_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                                          Inches(1), Inches(1.5), Inches(3.5), Inches(1.8))
    frontend_box.fill.solid()
    frontend_box.fill.fore_color.rgb = COLORS['light_green']
    frontend_box.line.color.rgb = COLORS['primary_green']
    add_body_text(slide, "前端层\nVue 3 + Element Plus\nECharts 数据可视化", 
                  Inches(1.2), Inches(1.8), Inches(3.1), Inches(1.2), 
                  font_size=14, color=COLORS['dark_green'], align=PP_ALIGN.CENTER)
    
    # 后端层
    backend_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                                         Inches(4.9), Inches(1.5), Inches(3.5), Inches(1.8))
    backend_box.fill.solid()
    backend_box.fill.fore_color.rgb = COLORS['accent_green']
    backend_box.line.color.rgb = COLORS['dark_green']
    add_body_text(slide, "后端层\nSpring Boot 3.3\nMyBatis-Plus + JWT", 
                  Inches(5.1), Inches(1.8), Inches(3.1), Inches(1.2), 
                  font_size=14, color=COLORS['white'], align=PP_ALIGN.CENTER)
    
    # 数据层
    db_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                                    Inches(8.8), Inches(1.5), Inches(3.5), Inches(1.8))
    db_box.fill.solid()
    db_box.fill.fore_color.rgb = COLORS['dark_green']
    db_box.line.color.rgb = COLORS['primary_green']
    add_body_text(slide, "数据层\nMySQL 9.1\n10张核心数据表", 
                  Inches(9), Inches(1.8), Inches(3.1), Inches(1.2), 
                  font_size=14, color=COLORS['white'], align=PP_ALIGN.CENTER)
    
    # 连接箭头
    arrow1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.6), Inches(2.2), Inches(1.2), Inches(0.4))
    arrow1.fill.solid()
    arrow1.fill.fore_color.rgb = COLORS['text_gray']
    
    arrow2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(7.5), Inches(2.2), Inches(1.2), Inches(0.4))
    arrow2.fill.solid()
    arrow2.fill.fore_color.rgb = COLORS['text_gray']
    
    # 技术特点
    add_title_shape(slide, "技术特点", margin, Inches(4), Inches(4), Inches(0.5), 
                    font_size=24, color=COLORS['primary_green'])
    
    features = [
        "前后端分离：RESTful API 通信",
        "响应式设计：适配多设备访问",
        "JWT认证：安全用户权限管理",
        "逻辑删除：MyBatis-Plus数据保护",
    ]
    add_bullet_points(slide, features, margin, Inches(4.6), Inches(12), Inches(2.5), 
                      font_size=16, color=COLORS['text_dark'])
    
    # ========== 第6页：技术栈详情 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg_color(slide, COLORS['bg_light'])
    
    add_title_shape(slide, "技术栈详情", margin, Inches(0.4), Inches(4), Inches(0.5), 
                    font_size=36, color=COLORS['dark_green'])
    add_colored_bar(slide, margin, Inches(0.95), Inches(1.2), Inches(0.05), COLORS['accent_green'])
    
    # 前端技术栈
    add_title_shape(slide, "前端技术栈", Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.4), 
                    font_size=22, color=COLORS['primary_green'])
    
    frontend_tech = [
        "Vue 3.5 - 渐进式JavaScript框架",
        "Element Plus 2.13 - UI组件库",
        "Vue Router 4.6 - 前端路由管理",
        "Axios 1.16 - HTTP请求库",
        "ECharts 6.0 - 数据可视化图表",
        "Vite 8.0 - 构建工具",
    ]
    add_bullet_points(slide, frontend_tech, Inches(0.8), Inches(1.8), Inches(5.5), Inches(2.8), 
                      font_size=15, color=COLORS['text_dark'])
    
    # 后端技术栈
    add_title_shape(slide, "后端技术栈", Inches(7), Inches(1.3), Inches(5.5), Inches(0.4), 
                    font_size=22, color=COLORS['primary_green'])
    
    backend_tech = [
        "Spring Boot 3.3.5 - 应用框架",
        "MyBatis-Plus 3.5.7 - ORM框架",
        "MySQL 9.1.0 - 关系型数据库",
        "JJWT 0.12.6 - JWT认证库",
        "Lombok - 代码简化工具",
        "Java 21 - 编程语言",
    ]
    add_bullet_points(slide, backend_tech, Inches(7), Inches(1.8), Inches(5.5), Inches(2.8), 
                      font_size=15, color=COLORS['text_dark'])
    
    # 底部部署信息
    add_colored_bar(slide, Inches(0.8), Inches(5.2), Inches(11.7), Inches(1.8), COLORS['primary_green'])
    
    deploy_info = [
        "前端端口：5173（Vite开发服务器）",
        "后端端口：8088（Spring Boot）",
        "数据库：MySQL 3306端口",
        "代理配置：Vite代理转发API请求",
    ]
    y_pos = Inches(5.4)
    for info in deploy_info:
        add_body_text(slide, "• " + info, Inches(1), y_pos, Inches(11), Inches(0.4), 
                      font_size=14, color=COLORS['white'])
        y_pos += Inches(0.38)
    
    # ========== 第7页：功能架构 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg_color(slide, COLORS['white'])
    
    add_title_shape(slide, "功能架构", margin, Inches(0.4), Inches(4), Inches(0.5), 
                    font_size=36, color=COLORS['dark_green'])
    add_colored_bar(slide, margin, Inches(0.95), Inches(1.2), Inches(0.05), COLORS['accent_green'])
    
    # 8大模块展示
    modules = [
        ("用户认证", "JWT登录\n权限管理"),
        ("家畜种类", "树形分类\n参数配置"),
        ("养殖场地", "场地管理\n容量配置"),
        ("存栏管理", "批次跟踪\n变动记录"),
        ("成本管理", "成本录入\n自动计算"),
        ("价格规则", "浮动策略\n条件定价"),
        ("收益预估", "智能计算\n利润分析"),
        ("养殖建议", "AI推荐\n模板管理"),
    ]
    
    positions = [
        (0.6, 1.4), (3.5, 1.4), (6.4, 1.4), (9.3, 1.4),
        (0.6, 3.9), (3.5, 3.9), (6.4, 3.9), (9.3, 3.9),
    ]
    
    colors = [COLORS['light_green'], COLORS['accent_green'], COLORS['primary_green'], COLORS['dark_green']]
    
    for i, ((title, desc), (x, y)) in enumerate(zip(modules, positions)):
        # 模块卡片
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                                     Inches(x), Inches(y), Inches(2.6), Inches(2.3))
        box.fill.solid()
        box.fill.fore_color.rgb = colors[i % 4]
        if i % 4 >= 2:
            text_color = COLORS['white']
        else:
            text_color = COLORS['dark_green']
        box.line.fill.background()
        
        # 标题
        add_body_text(slide, title, Inches(x + 0.15), Inches(y + 0.3), Inches(2.3), Inches(0.5), 
                      font_size=18, color=text_color, align=PP_ALIGN.CENTER)
        # 描述
        add_body_text(slide, desc, Inches(x + 0.15), Inches(y + 0.9), Inches(2.3), Inches(1.2), 
                      font_size=13, color=text_color, align=PP_ALIGN.CENTER)
    
    # ========== 第8页：数据库设计 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg_color(slide, COLORS['bg_light'])
    
    add_title_shape(slide, "数据库设计", margin, Inches(0.4), Inches(4), Inches(0.5), 
                    font_size=36, color=COLORS['dark_green'])
    add_colored_bar(slide, margin, Inches(0.95), Inches(1.2), Inches(0.05), COLORS['accent_green'])
    
    # 数据表列表
    tables = [
        ("sys_user", "系统用户表", "username, role"),
        ("livestock_category", "家畜种类表", "父子分类, 生长周期"),
        ("breeding_site", "养殖场地表", "容量, 位置信息"),
        ("batch", "存栏批次表", "批次号, 生长阶段"),
        ("batch_change", "存栏变动表", "变动类型, 数量"),
        ("cost_record", "成本记录表", "成本类型, 金额"),
        ("price_rule", "价格规则表", "浮动方向, 比例"),
        ("breeding_advice", "养殖建议表", "建议内容, 是否已读"),
        ("breeding_advice_template", "建议模板表", "种类, 阶段, 模板"),
        ("system_config", "系统配置表", "配置项, 值"),
    ]
    
    y_pos = Inches(1.4)
    for i, (table_name, table_desc, fields) in enumerate(tables):
        if i == 5:
            y_pos = Inches(1.4)
        
        x_offset = Inches(0.3) if i < 5 else Inches(6.8)
        
        # 表名
        add_title_shape(slide, table_name, x_offset, y_pos, Inches(3.5), Inches(0.35), 
                        font_size=14, color=COLORS['primary_green'])
        # 描述
        add_body_text(slide, f"{table_desc} | {fields}", x_offset, y_pos + Inches(0.35), 
                      Inches(6), Inches(0.3), font_size=12, color=COLORS['text_gray'])
        
        y_pos += Inches(0.75)
    
    # 设计特点
    add_colored_bar(slide, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.6), COLORS['primary_green'])
    
    features = [
        "逻辑删除：使用deleted字段保护数据",
        "自动填充：create_time/update_time自动填充",
        "关联设计：批次-场地-种类多表关联",
        "索引优化：常用查询字段建立索引",
    ]
    add_bullet_points(slide, features, Inches(0.7), Inches(5.7), Inches(11.8), Inches(1.3), 
                      font_size=15, color=COLORS['white'])
    
    # ========== 第9页：仪表盘 - 配截图 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg_color(slide, COLORS['white'])
    
    add_title_shape(slide, "仪表盘 - 数据可视化总览", margin, Inches(0.3), Inches(8), Inches(0.5), 
                    font_size=32, color=COLORS['dark_green'])
    add_colored_bar(slide, margin, Inches(0.85), Inches(2), Inches(0.04), COLORS['accent_green'])
    
    # 添加截图
    img_path = os.path.join(ASSETS_DIR, SCREENSHOTS[0][0])
    add_image_to_slide(slide, img_path, Inches(0.5), Inches(1.1), Inches(12.3), Inches(5.8))
    
    # 功能说明
    add_colored_bar(slide, Inches(0.5), Inches(7), Inches(12.3), Inches(0.35), COLORS['light_green'])
    add_body_text(slide, "核心指标：存栏总数、饲养批次、预估收入、预估利润 | 图表：品种分布、生长阶段、场地利用率",
                  Inches(0.7), Inches(7.02), Inches(12), Inches(0.3), 
                  font_size=14, color=COLORS['dark_green'])
    
    # ========== 第10页：基础数据管理 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg_color(slide, COLORS['bg_light'])
    
    add_title_shape(slide, "基础数据管理", margin, Inches(0.3), Inches(6), Inches(0.5), 
                    font_size=32, color=COLORS['dark_green'])
    add_colored_bar(slide, margin, Inches(0.85), Inches(1.8), Inches(0.04), COLORS['accent_green'])
    
    # 左侧：家畜种类
    add_title_shape(slide, "家畜种类管理", Inches(0.5), Inches(1.1), Inches(5.8), Inches(0.4), 
                    font_size=22, color=COLORS['primary_green'])
    
    img_path = os.path.join(ASSETS_DIR, SCREENSHOTS[1][0])
    add_image_to_slide(slide, img_path, Inches(0.5), Inches(1.6), Inches(5.8), Inches(3.2))
    
    add_body_text(slide, "• 树形分类结构：支持多级分类\n• 养殖参数配置：生长周期、饲料消耗\n• 基准价格设定：便于收益计算",
                  Inches(0.5), Inches(5), Inches(5.8), Inches(1.5), 
                  font_size=14, color=COLORS['text_dark'])
    
    # 右侧：养殖场地
    add_title_shape(slide, "养殖场地管理", Inches(7), Inches(1.1), Inches(5.8), Inches(0.4), 
                    font_size=22, color=COLORS['primary_green'])
    
    img_path = os.path.join(ASSETS_DIR, SCREENSHOTS[2][0])
    add_image_to_slide(slide, img_path, Inches(7), Inches(1.6), Inches(5.8), Inches(3.2))
    
    add_body_text(slide, "• 场地信息维护：编号、名称、位置\n• 容量配置：最大存栏数量设定\n• 场地利用率监控",
                  Inches(7), Inches(5), Inches(5.8), Inches(1.5), 
                  font_size=14, color=COLORS['text_dark'])
    
    # ========== 第11页：存栏生命周期管理 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg_color(slide, COLORS['white'])
    
    add_title_shape(slide, "存栏生命周期管理", margin, Inches(0.3), Inches(6), Inches(0.5), 
                    font_size=32, color=COLORS['dark_green'])
    add_colored_bar(slide, margin, Inches(0.85), Inches(2.2), Inches(0.04), COLORS['accent_green'])
    
    # 截图
    img_path = os.path.join(ASSETS_DIR, SCREENSHOTS[3][0])
    add_image_to_slide(slide, img_path, Inches(0.5), Inches(1.1), Inches(12.3), Inches(4.8))
    
    # 功能特点
    add_colored_bar(slide, Inches(0.5), Inches(6.1), Inches(12.3), Inches(1.1), COLORS['primary_green'])
    
    features = [
        "批次管理：自动生成唯一批次号，支持批次筛选与搜索",
        "生长阶段：实时跟踪批次当前生长阶段（苗种→育肥→成年→出栏前）",
        "变动记录：出栏、死亡、转群、补栏等变动全程记录",
    ]
    add_bullet_points(slide, features, Inches(0.7), Inches(6.25), Inches(11.8), Inches(0.9), 
                      font_size=14, color=COLORS['white'])
    
    # ========== 第12页：成本管理系统 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg_color(slide, COLORS['bg_light'])
    
    add_title_shape(slide, "成本管理系统 - 智能化计算", margin, Inches(0.3), Inches(7), Inches(0.5), 
                    font_size=32, color=COLORS['dark_green'])
    add_colored_bar(slide, margin, Inches(0.85), Inches(2.2), Inches(0.04), COLORS['accent_green'])
    
    # 截图
    img_path = os.path.join(ASSETS_DIR, SCREENSHOTS[4][0])
    add_image_to_slide(slide, img_path, Inches(6.5), Inches(1.2), Inches(6.3), Inches(4.5))
    
    # 左侧功能说明
    add_title_shape(slide, "成本类型", Inches(0.5), Inches(1.2), Inches(5.5), Inches(0.4), 
                    font_size=22, color=COLORS['primary_green'])
    
    cost_types = [
        "饲料成本：根据存栏天数自动计算",
        "苗种成本：购入时的初始投入",
        "防疫成本：疫苗、药品等费用",
        "人工成本：人工管理费用",
    ]
    add_bullet_points(slide, cost_types, Inches(0.5), Inches(1.7), Inches(5.5), Inches(2), 
                      font_size=16, color=COLORS['text_dark'])
    
    # 自动计算说明
    add_title_shape(slide, "智能计算", Inches(0.5), Inches(3.9), Inches(5.5), Inches(0.4), 
                    font_size=22, color=COLORS['primary_green'])
    
    auto_calc = [
        "饲料费 = 存栏数量 × 天数 × 日饲料消耗",
        "管理费 = 存栏数量 × 天数 × 日管理费",
        "支持批量自动计算，一键生成成本记录",
    ]
    add_bullet_points(slide, auto_calc, Inches(0.5), Inches(4.4), Inches(5.5), Inches(1.5), 
                      font_size=14, color=COLORS['text_dark'])
    
    # 底部统计
    add_colored_bar(slide, Inches(0.5), Inches(5.9), Inches(12.3), Inches(1.2), COLORS['accent_green'])
    add_body_text(slide, "成本记录支持按种类、成本类型、日期范围筛选查询，便于精细化成本分析",
                  Inches(0.7), Inches(6.2), Inches(12), Inches(0.6), 
                  font_size=16, color=COLORS['white'])
    
    # ========== 第13页：价格规则管理 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg_color(slide, COLORS['white'])
    
    add_title_shape(slide, "价格规则管理 - 灵活定价策略", margin, Inches(0.3), Inches(7), Inches(0.5), 
                    font_size=32, color=COLORS['dark_green'])
    add_colored_bar(slide, margin, Inches(0.85), Inches(2.4), Inches(0.04), COLORS['accent_green'])
    
    # 截图
    img_path = os.path.join(ASSETS_DIR, SCREENSHOTS[5][0])
    add_image_to_slide(slide, img_path, Inches(0.5), Inches(3.8), Inches(12.3), Inches(3.4))
    
    # 上方功能说明
    add_title_shape(slide, "规则配置", Inches(0.5), Inches(1.2), Inches(3.5), Inches(0.4), 
                    font_size=20, color=COLORS['primary_green'])
    
    rules = [
        "上浮/下浮：根据市场情况调整",
        "浮动比例：可设置具体百分比",
        "适用条件：按月份、季节等条件触发",
    ]
    add_bullet_points(slide, rules, Inches(0.5), Inches(1.65), Inches(3.5), Inches(1.8), 
                      font_size=14, color=COLORS['text_dark'])
    
    add_title_shape(slide, "应用场景", Inches(4.5), Inches(1.2), Inches(4), Inches(0.4), 
                    font_size=20, color=COLORS['primary_green'])
    
    scenarios = [
        "春节猪价上涨：+8%",
        "冬季火锅季羊肉上涨：+12%",
        "夏季家禽需求下降：-5%",
    ]
    add_bullet_points(slide, scenarios, Inches(4.5), Inches(1.65), Inches(4), Inches(1.8), 
                      font_size=14, color=COLORS['text_dark'])
    
    add_title_shape(slide, "优先级机制", Inches(9), Inches(1.2), Inches(3.5), Inches(0.4), 
                    font_size=20, color=COLORS['primary_green'])
    
    add_body_text(slide, "多条规则同时满足时，按优先级执行，确保定价策略的合理性",
                  Inches(9), Inches(1.65), Inches(3.5), Inches(1.8), 
                  font_size=14, color=COLORS['text_dark'])
    
    # ========== 第14页：收益预估系统（重点） ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg_color(slide, COLORS['bg_light'])
    
    add_title_shape(slide, "收益预估系统 - 智能决策支持", margin, Inches(0.3), Inches(7), Inches(0.5), 
                    font_size=32, color=COLORS['dark_green'])
    add_colored_bar(slide, margin, Inches(0.85), Inches(2.4), Inches(0.04), COLORS['accent_green'])
    
    # 截图
    img_path = os.path.join(ASSETS_DIR, SCREENSHOTS[6][0])
    add_image_to_slide(slide, img_path, Inches(0.5), Inches(2.8), Inches(12.3), Inches(3.8))
    
    # 顶部亮点展示
    add_colored_bar(slide, Inches(0.5), Inches(1.1), Inches(12.3), Inches(1.5), COLORS['primary_green'])
    
    highlights = [
        "预计总收入：¥8,470,164.64",
        "预估总成本：¥6,935,777.30",
        "预计总利润：¥1,534,387.34",
        "利润率：18.1%",
    ]
    
    x_pos = Inches(0.8)
    for highlight in highlights:
        add_body_text(slide, highlight, x_pos, Inches(1.25), Inches(3), Inches(0.5), 
                      font_size=15, color=COLORS['white'])
        x_pos += Inches(3.1)
    
    # 计算逻辑
    add_colored_bar(slide, Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.5), COLORS['accent_green'])
    add_body_text(slide, "计算逻辑：预计产肉量 × 调整后单价 = 预计收入 | 收入 - 总成本 = 利润",
                  Inches(0.7), Inches(6.85), Inches(12), Inches(0.4), 
                  font_size=14, color=COLORS['white'])
    
    # ========== 第15页：养殖建议系统 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg_color(slide, COLORS['white'])
    
    add_title_shape(slide, "AI养殖建议系统", margin, Inches(0.3), Inches(6), Inches(0.5), 
                    font_size=32, color=COLORS['dark_green'])
    add_colored_bar(slide, margin, Inches(0.85), Inches(1.8), Inches(0.04), COLORS['accent_green'])
    
    # 截图
    img_path = os.path.join(ASSETS_DIR, SCREENSHOTS[7][0])
    add_image_to_slide(slide, img_path, Inches(6), Inches(1.1), Inches(6.8), Inches(5.5))
    
    # 左侧功能
    add_title_shape(slide, "智能生成", Inches(0.5), Inches(1.3), Inches(5), Inches(0.4), 
                    font_size=22, color=COLORS['primary_green'])
    
    ai_features = [
        "根据家畜种类自动匹配",
        "按当前生长阶段推送",
        "基于存栏数量个性化",
        "支持自定义模板管理",
    ]
    add_bullet_points(slide, ai_features, Inches(0.5), Inches(1.8), Inches(5), Inches(1.8), 
                      font_size=16, color=COLORS['text_dark'])
    
    add_title_shape(slide, "建议类型", Inches(0.5), Inches(3.8), Inches(5), Inches(0.4), 
                    font_size=22, color=COLORS['primary_green'])
    
    advice_types = [
        "存栏预警：关注养殖密度",
        "出栏提醒：最佳出栏时机",
        "防疫建议：疫苗接种计划",
        "饲养指导：饲料配比优化",
    ]
    add_bullet_points(slide, advice_types, Inches(0.5), Inches(4.3), Inches(5), Inches(1.8), 
                      font_size=16, color=COLORS['text_dark'])
    
    # 底部强调
    add_colored_bar(slide, Inches(0.5), Inches(6.1), Inches(5.3), Inches(1), COLORS['light_green'])
    add_body_text(slide, "AI智能分析，\n让养殖更科学",
                  Inches(0.7), Inches(6.2), Inches(4.9), Inches(0.8), 
                  font_size=18, color=COLORS['dark_green'], align=PP_ALIGN.CENTER)
    
    # ========== 第16页：功能亮点总结 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg_color(slide, COLORS['bg_light'])
    
    add_title_shape(slide, "核心功能亮点", margin, Inches(0.4), Inches(5), Inches(0.5), 
                    font_size=36, color=COLORS['dark_green'])
    add_colored_bar(slide, margin, Inches(0.95), Inches(2), Inches(0.05), COLORS['accent_green'])
    
    # 4个核心亮点
    highlights = [
        ("01", "全生命周期管理", "从入场到出栏全程跟踪\n批次管理与变动记录"),
        ("02", "智能成本计算", "自动计算饲料费、管理费\n精细化成本分析"),
        ("03", "收益预估决策", "基于价格规则智能计算\n辅助出栏决策"),
        ("04", "AI养殖建议", "按生长阶段推送建议\n模板化建议管理"),
    ]
    
    positions = [
        (1, 1.6), (7, 1.6),
        (1, 4.3), (7, 4.3),
    ]
    
    for (num, title, desc), (x, y) in zip(highlights, positions):
        # 卡片背景
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                                      Inches(x), Inches(y), Inches(5.3), Inches(2.3))
        card.fill.solid()
        card.fill.fore_color.rgb = COLORS['white']
        card.line.color.rgb = COLORS['light_green']
        
        # 编号
        num_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, 
                                            Inches(x + 0.3), Inches(y + 0.3), Inches(0.7), Inches(0.7))
        num_circle.fill.solid()
        num_circle.fill.fore_color.rgb = COLORS['accent_green']
        num_circle.line.fill.background()
        
        add_body_text(slide, num, Inches(x + 0.35), Inches(y + 0.4), Inches(0.6), Inches(0.5), 
                      font_size=20, color=COLORS['white'], align=PP_ALIGN.CENTER)
        
        # 标题
        add_title_shape(slide, title, Inches(x + 1.2), Inches(y + 0.35), Inches(3.8), Inches(0.5), 
                        font_size=22, color=COLORS['primary_green'])
        
        # 描述
        add_body_text(slide, desc, Inches(x + 0.3), Inches(y + 1.1), Inches(4.7), Inches(1), 
                      font_size=15, color=COLORS['text_dark'])
    
    # ========== 第17页：数据可视化 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg_color(slide, COLORS['white'])
    
    add_title_shape(slide, "数据可视化展示", margin, Inches(0.4), Inches(5), Inches(0.5), 
                    font_size=36, color=COLORS['dark_green'])
    add_colored_bar(slide, margin, Inches(0.95), Inches(2), Inches(0.05), COLORS['accent_green'])
    
    # 图表类型展示
    charts = [
        ("仪表盘图表", "品种存栏分布", "环形图展示各类家畜占比"),
        ("柱状图", "生长阶段分布", "直观显示各阶段数量"),
        ("进度条", "场地利用率", "实时监控场地使用情况"),
        ("数据卡片", "关键指标", "收入、利润等核心数据"),
    ]
    
    y_pos = Inches(1.5)
    for title, chart_name, desc in charts:
        # 标题
        add_title_shape(slide, title, Inches(0.8), y_pos, Inches(3), Inches(0.4), 
                        font_size=18, color=COLORS['primary_green'])
        
        # 图表名
        add_body_text(slide, chart_name, Inches(4), y_pos, Inches(3), Inches(0.4), 
                      font_size=16, color=COLORS['text_dark'])
        
        # 描述
        add_body_text(slide, desc, Inches(7.5), y_pos, Inches(5), Inches(0.4), 
                      font_size=14, color=COLORS['text_gray'])
        
        y_pos += Inches(0.9)
    
    # 技术实现
    add_colored_bar(slide, Inches(0.8), Inches(5.2), Inches(11.7), Inches(1.8), COLORS['primary_green'])
    
    add_title_shape(slide, "技术实现：ECharts 6.0 + Vue-ECharts", Inches(1), Inches(5.4), Inches(10), Inches(0.4), 
                    font_size=18, color=COLORS['white'])
    
    tech_desc = [
        "响应式设计：图表自动适应容器大小",
        "主题定制：匹配系统绿色系配色方案",
        "交互动画：悬停提示、渐变效果",
    ]
    add_bullet_points(slide, tech_desc, Inches(1), Inches(5.85), Inches(10), Inches(1), 
                      font_size=14, color=COLORS['white'])
    
    # ========== 第18页：应用价值 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg_color(slide, COLORS['bg_light'])
    
    add_title_shape(slide, "应用价值", margin, Inches(0.4), Inches(4), Inches(0.5), 
                    font_size=36, color=COLORS['dark_green'])
    add_colored_bar(slide, margin, Inches(0.95), Inches(1.2), Inches(0.05), COLORS['accent_green'])
    
    # 3大价值
    values = [
        ("效率提升", "数字化管理", "• 数据录入效率提升80%\n• 查询统计一键完成\n• 告别纸质记录"),
        ("科学决策", "数据驱动", "• 收益预估辅助出栏决策\n• 成本分析优化投入\n• 价格规则把握市场"),
        ("风险管控", "全程追溯", "• 批次变动全程记录\n• 建议系统降低风险\n• 数据备份防丢失"),
    ]
    
    positions = [(0.7, 1.5), (4.6, 1.5), (8.5, 1.5)]
    
    for (title, subtitle, content), (x, y) in zip(values, positions):
        # 卡片
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                                      Inches(x), Inches(y), Inches(3.8), Inches(4.5))
        card.fill.solid()
        card.fill.fore_color.rgb = COLORS['white']
        card.line.color.rgb = COLORS['light_green']
        
        # 标题
        add_title_shape(slide, title, Inches(x + 0.2), Inches(y + 0.3), Inches(3.4), Inches(0.5), 
                        font_size=24, color=COLORS['primary_green'])
        
        # 副标题
        add_body_text(slide, subtitle, Inches(x + 0.2), Inches(y + 0.85), Inches(3.4), Inches(0.4), 
                      font_size=16, color=COLORS['text_gray'])
        
        # 内容
        add_body_text(slide, content, Inches(x + 0.2), Inches(y + 1.4), Inches(3.4), Inches(2.8), 
                      font_size=14, color=COLORS['text_dark'])
    
    # 底部数据
    add_colored_bar(slide, Inches(0.7), Inches(6.2), Inches(11.7), Inches(0.9), COLORS['primary_green'])
    add_body_text(slide, "系统已配置：8个家畜种类 | 4个养殖场地 | 14个存栏批次 | 3000+头存栏",
                  Inches(1), Inches(6.4), Inches(11), Inches(0.5), 
                  font_size=18, color=COLORS['white'])
    
    # ========== 第19页：创新点总结 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg_color(slide, COLORS['white'])
    
    add_title_shape(slide, "项目创新点", margin, Inches(0.4), Inches(4), Inches(0.5), 
                    font_size=36, color=COLORS['dark_green'])
    add_colored_bar(slide, margin, Inches(0.95), Inches(1.5), Inches(0.05), COLORS['accent_green'])
    
    innovations = [
        "业务创新：覆盖养殖全生命周期的数字化管理",
        "技术创新：Vue3 + Spring Boot 现代化技术栈",
        "算法创新：智能成本计算与收益预估模型",
        "交互创新：ECharts数据可视化大屏展示",
        "设计创新：清新绿色系界面，符合农业主题",
    ]
    
    y_pos = Inches(1.6)
    for i, innovation in enumerate(innovations, 1):
        # 编号圆圈
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, 
                                        Inches(0.8), y_pos, Inches(0.6), Inches(0.6))
        circle.fill.solid()
        circle.fill.fore_color.rgb = COLORS['accent_green']
        circle.line.fill.background()
        
        add_body_text(slide, str(i), Inches(0.85), y_pos + Inches(0.05), Inches(0.5), Inches(0.5), 
                      font_size=18, color=COLORS['white'], align=PP_ALIGN.CENTER)
        
        # 内容
        add_body_text(slide, innovation, Inches(1.6), y_pos + Inches(0.05), Inches(10.5), Inches(0.5), 
                      font_size=18, color=COLORS['text_dark'])
        
        y_pos += Inches(0.9)
    
    # 底部对比
    add_colored_bar(slide, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.9), COLORS['primary_green'])
    add_body_text(slide, "从传统手工记录到数字化智能管理，本项目为智慧农业发展提供了实践参考",
                  Inches(1), Inches(6.4), Inches(11), Inches(0.5), 
                  font_size=16, color=COLORS['white'])
    
    # ========== 第20页：项目总结 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg_color(slide, COLORS['bg_light'])
    
    add_title_shape(slide, "项目总结", margin, Inches(0.4), Inches(4), Inches(0.5), 
                    font_size=36, color=COLORS['dark_green'])
    add_colored_bar(slide, margin, Inches(0.95), Inches(1.5), Inches(0.05), COLORS['accent_green'])
    
    # 左侧：完成内容
    add_title_shape(slide, "已完成", Inches(0.8), Inches(1.4), Inches(5), Inches(0.4), 
                    font_size=24, color=COLORS['primary_green'])
    
    completed = [
        "8大功能模块全部实现",
        "前后端分离架构搭建",
        "10张数据库表设计",
        "8张系统界面截图",
        "数据可视化图表集成",
    ]
    add_bullet_points(slide, completed, Inches(0.8), Inches(1.9), Inches(5), Inches(2.5), 
                      font_size=16, color=COLORS['text_dark'])
    
    # 右侧：技术收获
    add_title_shape(slide, "技术收获", Inches(7), Inches(1.4), Inches(5), Inches(0.4), 
                    font_size=24, color=COLORS['primary_green'])
    
    gains = [
        "Vue3组合式API实践",
        "Spring Boot RESTful开发",
        "MyBatis-Plus高效CRUD",
        "JWT认证与权限控制",
        "ECharts数据可视化",
    ]
    add_bullet_points(slide, gains, Inches(7), Inches(1.9), Inches(5), Inches(2.5), 
                      font_size=16, color=COLORS['text_dark'])
    
    # 底部项目信息
    add_colored_bar(slide, Inches(0.8), Inches(5), Inches(11.7), Inches(2.1), COLORS['primary_green'])
    
    add_title_shape(slide, "项目信息", Inches(1), Inches(5.2), Inches(5), Inches(0.4), 
                    font_size=20, color=COLORS['white'])
    
    info_items = [
        "项目名称：家畜养殖信息采集及收益估计系统",
        "技术栈：Vue3 + Spring Boot + MySQL",
        "代码量：约8000行（前端5000 + 后端3000）",
    ]
    add_bullet_points(slide, info_items, Inches(1), Inches(5.65), Inches(10), Inches(1.3), 
                      font_size=14, color=COLORS['white'])
    
    # ========== 第21页：致谢 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg_color(slide, COLORS['primary_green'])
    
    # 大标题
    title_shape = slide.shapes.add_textbox(Inches(0), Inches(2.5), slide_width, Inches(1))
    tf = title_shape.text_frame
    p = tf.paragraphs[0]
    p.text = "感谢聆听"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    subtitle_shape = slide.shapes.add_textbox(Inches(0), Inches(3.8), slide_width, Inches(0.6))
    tf2 = subtitle_shape.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "敬请批评指正"
    p2.font.size = Pt(32)
    p2.font.color.rgb = COLORS['light_green']
    p2.alignment = PP_ALIGN.CENTER
    
    # Q&A
    qa_shape = slide.shapes.add_textbox(Inches(0), Inches(5.5), slide_width, Inches(0.5))
    tf3 = qa_shape.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "Q & A"
    p3.font.size = Pt(28)
    p3.font.color.rgb = COLORS['white']
    p3.alignment = PP_ALIGN.CENTER
    
    # 装饰圆
    for i in range(3):
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, 
                                        Inches(2 + i * 3.5), Inches(6.2), Inches(0.5), Inches(0.5))
        circle.fill.solid()
        circle.fill.fore_color.rgb = COLORS['light_green']
        circle.line.fill.background()
    
    # 保存文件
    prs.save(OUTPUT_FILE)
    print(f"PPT已保存到: {OUTPUT_FILE}")
    print(f"总页数: {len(prs.slides)}")
    
    return OUTPUT_FILE

if __name__ == "__main__":
    create_presentation()
